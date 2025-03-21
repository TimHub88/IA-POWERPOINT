import os
import tempfile
from typing import Optional, Tuple, Dict
import uuid
import httpx
from io import BytesIO
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt

from ..domain.entities import Presentation, Slide
from ..domain.repository import PresentationRepository
from .pexels_client import PexelsClient


class PPTXGenerator(PresentationRepository):
    # Image de fallback si Pexels échoue
    FALLBACK_IMAGE = "/static/images/fallback.jpg"
    FALLBACK_IMAGE_PATH = "static/images/fallback.jpg"
    
    def __init__(self, output_dir: str = "static/presentations"):
        self.output_dir = output_dir
        self.pexels_client = PexelsClient()
        
        # Ensure output directory exists
        os.makedirs(self.output_dir, exist_ok=True)
        print(f"🚀 PPTXGenerator initialized with output directory: {self.output_dir}")
        
        # Vérifier que l'image de fallback existe
        if not os.path.exists(self.FALLBACK_IMAGE_PATH):
            print(f"⚠️ Fallback image not found at {self.FALLBACK_IMAGE_PATH}. Will use placeholder URLs.")
    
    async def save(self, presentation: Presentation, filename: str) -> str:
        """
        Save the presentation to a PowerPoint file.
        
        Args:
            presentation: The presentation object
            filename: The filename to save as
            
        Returns:
            The path to the saved file
        """
        print(f"🛠️ Starting to create PowerPoint file: {filename}")
        print(f"📊 Presentation contains {len(presentation.slides)} slides")
        
        # Create a new PowerPoint presentation
        pptx = PPTXPresentation()
        
        # Remove the default slide
        if len(pptx.slides) > 0:
            print("🔄 Removing default slide from template")
            r_id = pptx.slides._sldIdLst[0].rId
            pptx.part.drop_rel(r_id)
            pptx.slides._sldIdLst.remove(pptx.slides._sldIdLst[0])
        
        # Add slides
        print("🔄 Creating HTTP client for image downloads")
        async with httpx.AsyncClient(timeout=30.0) as client:
            for i, slide in enumerate(presentation.slides):
                print(f"📑 Processing slide {i+1}/{len(presentation.slides)}: '{slide.title}'")
                await self._add_slide(pptx, slide, client)
        
        # Save the presentation
        file_path = os.path.join(self.output_dir, filename)
        print(f"💾 Saving PowerPoint file to: {file_path}")
        pptx.save(file_path)
        
        print(f"✅ PowerPoint file saved successfully")
        # Return the relative path to be used in URLs
        return os.path.join("presentations", filename)
    
    async def _add_slide(self, pptx: PPTXPresentation, slide: Slide, client: httpx.AsyncClient) -> None:
        """
        Add a slide to the PowerPoint presentation.
        
        Args:
            pptx: The PowerPoint presentation object
            slide: The slide to add
            client: HTTPx client for downloading images
        """
        # Add a slide with a title and content layout
        print(f"🔄 Adding new slide with title: '{slide.title}'")
        layout = pptx.slide_layouts[1]  # Title and Content layout
        pptx_slide = pptx.slides.add_slide(layout)
        
        # Set the title
        title = pptx_slide.shapes.title
        title.text = slide.title
        print(f"✍️ Added slide title: '{slide.title}'")
        
        # Set the content
        content = pptx_slide.placeholders[1]
        content.text = slide.description
        desc_preview = slide.description[:50] + "..." if len(slide.description) > 50 else slide.description
        print(f"📝 Added slide content: '{desc_preview}'")
        
        # Format text (optional)
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
        print("🎨 Applied text formatting")
        
        # Vérifier si le mode sans images est activé (pas de keywords)
        if not slide.keywords:
            print(f"ℹ️ Skip image processing - images disabled for slide: {slide.title}")
            return
        
        # Get relevant image for the slide based on keywords
        print(f"🖼️ Processing image for slide: {slide.title}")
        
        # Search for a relevant image using keywords
        image_url = await self._get_image_for_slide(slide)
        print(f"🔗 Using image URL: {image_url}")
        
        try:
            # Check if it's a local file path (starting with /static)
            if image_url.startswith("/static"):
                local_path = image_url[1:]  # Remove leading slash
                if os.path.exists(local_path):
                    print(f"✅ Using local image file: {local_path}")
                    image_data = open(local_path, 'rb').read()
                    image_ext = os.path.splitext(local_path)[1][1:]  # Get extension without dot
                    self._add_image_to_slide(pptx_slide, image_data, image_ext, slide.title)
                else:
                    print(f"⚠️ Local image file not found: {local_path}")
            else:
                # Download the image from URL
                print(f"📥 Downloading image from: {image_url}")
                image_data, image_ext = await self._download_image(image_url, client)
                
                if image_data:
                    print(f"✅ Image downloaded successfully ({len(image_data)} bytes, format: {image_ext})")
                    self._add_image_to_slide(pptx_slide, image_data, image_ext, slide.title)
                else:
                    print(f"⚠️ Failed to download image for slide: {slide.title}")
                    # Try fallback local image
                    if os.path.exists(self.FALLBACK_IMAGE_PATH):
                        print(f"🔄 Using local fallback image: {self.FALLBACK_IMAGE_PATH}")
                        image_data = open(self.FALLBACK_IMAGE_PATH, 'rb').read()
                        image_ext = os.path.splitext(self.FALLBACK_IMAGE_PATH)[1][1:]
                        self._add_image_to_slide(pptx_slide, image_data, image_ext, slide.title)
        except Exception as e:
            print(f"⚠️ Error adding image to slide: {str(e)}")
            import traceback
            print(f"⚠️ Traceback: {traceback.format_exc()}")
            # Continue without the image if there's an error
    
    def _add_image_to_slide(self, pptx_slide, image_data, image_ext, slide_title):
        """
        Add an image to a slide with proper centering.
        
        Args:
            pptx_slide: The PowerPoint slide object
            image_data: Binary image data
            image_ext: Image file extension
            slide_title: Title of the slide (for logging)
        """
        try:
            # Create a temporary file to save the image
            with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{image_ext}') as temp_file:
                temp_file.write(image_data)
                temp_file_path = temp_file.name
                print(f"💾 Image saved to temporary file: {temp_file_path}")
            
            try:
                # Get slide dimensions (standard PowerPoint is 10" x 7.5")
                slide_width = Inches(10)
                slide_height = Inches(7.5)
                
                # Set image dimensions preserving aspect ratio
                image_width = Inches(7)  # Larger width for better visibility
                
                # Calculate center position
                left = (slide_width - image_width) / 2
                
                # Position below the content - fixed position for better alignment
                # This places the image in the lower part of the slide, leaving room for text
                top = Inches(3.5)
                
                # Add the image to the slide
                print(f"🛠️ Adding image to slide with centered positioning")
                pptx_slide.shapes.add_picture(
                    temp_file_path, 
                    left, 
                    top, 
                    width=image_width
                )
                print(f"✅ Successfully added image to slide: {slide_title}")
            finally:
                # Clean up the temporary file
                print(f"🧹 Cleaning up temporary file")
                os.unlink(temp_file_path)
        except Exception as e:
            print(f"⚠️ Error in _add_image_to_slide: {str(e)}")
            import traceback
            print(f"⚠️ Traceback: {traceback.format_exc()}")
    
    async def _get_image_for_slide(self, slide: Slide) -> str:
        """
        Get an image URL for a slide based on its keywords.
        
        Args:
            slide: The slide to get an image for
            
        Returns:
            URL of an image for the slide
        """
        # If slide already has an image URL that's not from Picsum, use it
        if slide.image and "picsum.photos" not in slide.image and slide.image.startswith("http"):
            print(f"ℹ️ Slide already has a valid image URL: {slide.image}")
            return slide.image
            
        # If we have keywords, search using Pexels
        if slide.keywords:
            print(f"🔍 Using keywords for image search: {', '.join(slide.keywords)}")
            # Search for image using keywords
            return await self.pexels_client.search_image(
                keywords=slide.keywords,
                fallback_url=self.FALLBACK_IMAGE
            )
        
        # If no keywords but we have title/description, use those
        if not slide.keywords and (slide.title or slide.description):
            # Generate keywords from title and first sentence of description
            title_words = slide.title.split()
            if slide.description:
                first_sentence = slide.description.split('.')[0]
                desc_words = first_sentence.split()
            else:
                desc_words = []
                
            # Combine unique words from title and description
            combined_words = list(set(title_words + desc_words))
            
            # Filter out common words and limit to 5 words
            stop_words = {'a', 'an', 'the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with'}
            keywords = [word for word in combined_words if word.lower() not in stop_words][:5]
            
            if keywords:
                print(f"🔍 Generated keywords from content: {', '.join(keywords)}")
                return await self.pexels_client.search_image(
                    keywords=keywords,
                    fallback_url=self.FALLBACK_IMAGE
                )
        
        # If all else fails, return fallback image
        print(f"ℹ️ No keywords available, using fallback image")
        return self.FALLBACK_IMAGE
    
    async def _download_image(self, image_url: str, client: httpx.AsyncClient) -> Tuple[Optional[bytes], str]:
        """
        Download an image from a URL.
        
        Args:
            image_url: URL of the image to download
            client: HTTPx client
            
        Returns:
            Tuple of (image_data, image_extension) or (None, '') if download failed
        """
        try:
            print(f"📥 Attempting to download image from: {image_url}")
            response = await client.get(image_url, follow_redirects=True, timeout=20.0)
            print(f"🔄 Got response with status code: {response.status_code}")
            
            response.raise_for_status()
            
            # Get the content type to determine the image extension
            content_type = response.headers.get('content-type', '')
            ext = self._get_extension_from_content_type(content_type)
            
            # Log extra information about the response
            print(f"ℹ️ Response headers: {dict(response.headers)}")
            print(f"✅ Successfully downloaded image ({len(response.content)} bytes, type: {content_type}, extension: {ext})")
            return response.content, ext
        except Exception as e:
            print(f"⚠️ Error downloading image from {image_url}: {str(e)}")
            import traceback
            print(f"⚠️ Traceback: {traceback.format_exc()}")
            return None, ''
    
    def _get_extension_from_content_type(self, content_type: str) -> str:
        """
        Get the file extension from the content type.
        
        Args:
            content_type: Content type from the HTTP response
            
        Returns:
            File extension (without the dot)
        """
        print(f"🔍 Determining file extension from content type: '{content_type}'")
        if 'jpeg' in content_type or 'jpg' in content_type:
            ext = 'jpg'
        elif 'png' in content_type:
            ext = 'png'
        elif 'gif' in content_type:
            ext = 'gif'
        elif 'bmp' in content_type:
            ext = 'bmp'
        elif 'webp' in content_type:
            ext = 'webp'
        else:
            # Default to jpg if we can't determine the type
            ext = 'jpg'
        
        print(f"✅ Determined file extension: '{ext}'")
        return ext 
